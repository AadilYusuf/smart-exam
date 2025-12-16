import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text(file_path, file_type):
    if file_type == "pdf":
        import pdfplumber
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text()
        return text
    elif file_type == "docx":
        from docx import Document
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif file_type == "pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file_type in ["png", "jpg", "jpeg"]:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(file_path))
